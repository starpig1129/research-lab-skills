import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

from svg_to_pptx.converter import CoordSystem
from svg_to_pptx.connector import dispatch_connector, build_anchor_map, bind_connector_end, _add_line
from svg_to_pptx.style_parser import compute_style

CS = CoordSystem(svg_w=1200.0, svg_h=675.0)


def _slide():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_line_creates_connector():
    slide = _slide()
    elem = etree.fromstring('<line x1="100" y1="100" x2="400" y2="100" stroke="#000"/>')
    dispatch_connector(slide, elem, compute_style(elem, {}), CS)
    assert len(slide.shapes) == 1


def test_polyline_creates_n_minus_1_connectors():
    slide = _slide()
    elem = etree.fromstring('<polyline points="0,0 100,0 200,100" stroke="#000"/>')
    dispatch_connector(slide, elem, compute_style(elem, {}), CS)
    assert len(slide.shapes) == 2


def test_polygon_closes_path():
    slide = _slide()
    elem = etree.fromstring('<polygon points="0,0 100,0 50,100" stroke="#000"/>')
    dispatch_connector(slide, elem, compute_style(elem, {}), CS)
    assert len(slide.shapes) == 3


def test_build_anchor_map_rect():
    anchors = build_anchor_map([(None, 100.0, 50.0, 200.0, 100.0, 42)])
    assert 42 in anchors
    assert len(anchors[42]) == 8


def test_build_anchor_map_nearest():
    anchors = build_anchor_map([(None, 100.0, 50.0, 200.0, 100.0, 42)])
    top_center = next(a for a in anchors[42] if a[2] == 1)
    assert top_center[0] == 200.0
    assert top_center[1] == 50.0


def test_bind_connector_injects_xml():
    slide = _slide()
    conn = _add_line(slide, 0, 0, 100, 0, {})
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(CS.x(200)), Emu(CS.y(100)))
    bind_connector_end(conn, True, shape.shape_id, 3)
    cxnSp = conn._element
    nvCxnSpPr = cxnSp.find(qn('p:nvCxnSpPr'))
    cNvCxnSpPr = nvCxnSpPr.find(qn('p:cNvCxnSpPr'))
    stCxn = cNvCxnSpPr.find(qn('a:stCxn'))
    assert stCxn is not None
    assert stCxn.get('id') == str(shape.shape_id)
    assert stCxn.get('idx') == '3'
