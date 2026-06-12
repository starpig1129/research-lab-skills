"""test_integration.py — end-to-end SVG → PPTX shape verification."""
import io, os, tempfile
import pytest
from pptx import Presentation
from pptx.util import Emu

from svg_to_pptx.converter import SvgConverter, convert_file

_DIAGRAM_SVG = """\
<svg viewBox="0 0 1200 675" xmlns="http://www.w3.org/2000/svg">
  <rect x="40" y="80" width="160" height="70" fill="#3b82f6" stroke="#1d4ed8" stroke-width="2"/>
  <text x="120" y="112" fill="white" text-anchor="middle" font-size="14" font-weight="bold">Deep Research</text>
  <text x="120" y="132" fill="white" text-anchor="middle" font-size="11">13-agent pipeline</text>
  <circle cx="600" cy="337" r="50" fill="#10b981"/>
  <text x="600" y="341" fill="white" text-anchor="middle" font-size="13">Hub</text>
  <line x1="200" y1="115" x2="550" y2="337" stroke="#6b7280" stroke-width="1.5"/>
</svg>"""

_BEZIER_SVG = """\
<svg viewBox="0 0 400 200" xmlns="http://www.w3.org/2000/svg">
  <path d="M 0 100 C 100 0 300 0 400 100" stroke="#3b82f6" fill="none" stroke-width="3"/>
</svg>"""

_ARC_SVG = """\
<svg viewBox="0 0 400 400" xmlns="http://www.w3.org/2000/svg">
  <path d="M 200 50 A 150 150 0 0 1 350 200" stroke="#f59e0b" fill="none" stroke-width="2"/>
</svg>"""


def _make_tmp_svg(content: str) -> str:
    f = tempfile.NamedTemporaryFile(suffix=".svg", mode="w", delete=False)
    f.write(content); f.close()
    return f.name


def _conv(svg_content: str):
    path = _make_tmp_svg(svg_content)
    try:
        prs = Presentation()
        prs.slide_width = Emu(12_192_000)
        prs.slide_height = Emu(6_858_000)
        conv = SvgConverter(path)
        slide = conv.convert(prs, prs.slide_layouts[6])
        return slide, prs
    finally:
        os.unlink(path)


def test_diagram_shape_count():
    slide, _ = _conv(_DIAGRAM_SVG)
    assert len(slide.shapes) == 3


def test_diagram_text_in_rect():
    slide, _ = _conv(_DIAGRAM_SVG)
    rect = slide.shapes[0]
    assert rect.has_text_frame
    texts = [r.text for p in rect.text_frame.paragraphs for r in p.runs]
    assert any("Deep Research" in t for t in texts)


def test_diagram_connector_present():
    slide, _ = _conv(_DIAGRAM_SVG)
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    connectors = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(connectors) >= 1


def test_bezier_path_saves_valid_pptx():
    slide, prs = _conv(_BEZIER_SVG)
    assert len(slide.shapes) == 1
    buf = io.BytesIO()
    prs.save(buf)
    assert len(buf.getvalue()) > 5000


def test_arc_path_saves_valid_pptx():
    slide, prs = _conv(_ARC_SVG)
    assert len(slide.shapes) == 1
    buf = io.BytesIO()
    prs.save(buf)
    assert len(buf.getvalue()) > 5000


def test_convert_file_creates_pptx(tmp_path):
    svg_dir = tmp_path / "slides"
    svg_dir.mkdir()
    (svg_dir / "slide01.svg").write_text(_DIAGRAM_SVG)
    out = str(tmp_path / "deck.pptx")
    convert_file(str(svg_dir), out)
    assert os.path.exists(out)
    prs = Presentation(out)
    assert len(prs.slides) == 1


def test_all_shapes_have_positive_size():
    slide, _ = _conv(_DIAGRAM_SVG)
    for shape in slide.shapes:
        assert shape.width > 0
        assert shape.height > 0
