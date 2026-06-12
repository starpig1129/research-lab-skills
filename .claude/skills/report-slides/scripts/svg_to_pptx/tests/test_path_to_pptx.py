import io
import pytest
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

from svg_to_pptx.converter import CoordSystem
from svg_to_pptx.path_parser import parse_path
from svg_to_pptx.path_to_pptx import add_path_shape, arc_to_beziers

CS = CoordSystem(svg_w=1200.0, svg_h=675.0)
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _slide():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6]), prs


def test_straight_path_creates_shape():
    slide, _ = _slide()
    cmds = parse_path("M 100 100 L 300 100 L 300 200 Z")
    shape = add_path_shape(slide, cmds, CS, {"fill": "#3b82f6"})
    assert shape is not None
    assert len(slide.shapes) == 1


def test_straight_path_uses_custgeom():
    slide, _ = _slide()
    cmds = parse_path("M 100 100 L 300 100 L 300 200 Z")
    shape = add_path_shape(slide, cmds, CS, {})
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    custGeom = spPr.find(f"{{{_A_NS}}}custGeom")
    assert custGeom is not None


def test_bezier_path_injects_cubicBezTo():
    slide, _ = _slide()
    cmds = parse_path("M 0 100 C 25 0 75 0 100 100")
    shape = add_path_shape(slide, cmds, CS, {})
    sp = shape._element
    spPr = sp.find(qn("p:spPr"))
    custGeom = spPr.find(f"{{{_A_NS}}}custGeom")
    pathLst = custGeom.find(f"{{{_A_NS}}}pathLst")
    path = pathLst.find(f"{{{_A_NS}}}path")
    bezTo = path.find(f"{{{_A_NS}}}cubicBezTo")
    assert bezTo is not None
    assert len(bezTo) == 3


def test_arc_to_beziers_returns_list():
    result = arc_to_beziers(100, 50, 50, 50, 0, 0, 1, 0, 50)
    assert isinstance(result, list)
    assert len(result) >= 1
    for seg in result:
        assert len(seg) == 3


def test_arc_path_saves_without_error():
    slide, prs = _slide()
    cmds = parse_path("M 600 100 A 100 100 0 0 1 800 100")
    shape = add_path_shape(slide, cmds, CS, {"stroke": "#000000"})
    assert shape is not None
    buf = io.BytesIO()
    prs.save(buf)
    assert len(buf.getvalue()) > 1000


def test_empty_commands_returns_none():
    slide, _ = _slide()
    assert add_path_shape(slide, [], CS, {}) is None
