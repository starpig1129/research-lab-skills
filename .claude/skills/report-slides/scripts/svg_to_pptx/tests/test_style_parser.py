import pytest
from svg_to_pptx.converter import CoordSystem

def test_coordsystem_from_viewbox():
    cs = CoordSystem.from_viewbox("0 0 1200 675")
    assert cs.svg_w == 1200.0
    assert cs.svg_h == 675.0

def test_coordsystem_x_scaling():
    cs = CoordSystem.from_viewbox("0 0 1200 675")
    assert cs.x(600) == 6_096_000

def test_coordsystem_y_scaling():
    cs = CoordSystem.from_viewbox("0 0 1200 675")
    assert cs.y(337.5) == 3_429_000

def test_coordsystem_missing_viewbox():
    cs = CoordSystem.from_viewbox("")
    assert cs.svg_w == 1200.0
    assert cs.svg_h == 675.0

def test_coordsystem_malformed_viewbox():
    cs = CoordSystem.from_viewbox("0 0 abc 675")
    assert cs.svg_w == 1200.0
    assert cs.svg_h == 675.0


from svg_to_pptx.style_parser import (
    parse_inline_style, resolve_color, compute_style, apply_fill, apply_stroke
)
from pptx.dml.color import RGBColor
from lxml import etree


def test_parse_inline_style():
    r = parse_inline_style("fill:#3b82f6;stroke:#1d4ed8;stroke-width:2")
    assert r == {"fill": "#3b82f6", "stroke": "#1d4ed8", "stroke-width": "2"}


def test_parse_inline_style_with_spaces():
    r = parse_inline_style(" fill : #fff ; stroke: none ")
    assert r["fill"] == "#fff"
    assert r["stroke"] == "none"


def test_resolve_color_hex6():
    assert resolve_color("#3b82f6") == RGBColor(0x3b, 0x82, 0xf6)


def test_resolve_color_hex3():
    assert resolve_color("#f0f") == RGBColor(0xff, 0x00, 0xff)


def test_resolve_color_named_white():
    assert resolve_color("white") == RGBColor(0xff, 0xff, 0xff)


def test_resolve_color_named_blue():
    assert resolve_color("blue") == RGBColor(0x00, 0x00, 0xff)


def test_resolve_color_none():
    assert resolve_color("none") is None


def test_resolve_color_invalid():
    assert resolve_color("notacolor") is None


def test_compute_style_attr_fallback():
    elem = etree.fromstring('<rect fill="#ff0000"/>')
    result = compute_style(elem, {})
    assert result["fill"] == "#ff0000"


def test_compute_style_inherited():
    elem = etree.fromstring('<rect/>')
    result = compute_style(elem, {"fill": "#00ff00"})
    assert result["fill"] == "#00ff00"


def test_compute_style_inline_overrides_attr():
    elem = etree.fromstring('<rect fill="#ff0000" style="fill:#00ff00"/>')
    result = compute_style(elem, {})
    assert result["fill"] == "#00ff00"


def test_compute_style_inline_overrides_inherited():
    elem = etree.fromstring('<rect style="stroke:none"/>')
    result = compute_style(elem, {"stroke": "#000000"})
    assert result["stroke"] == "none"


from svg_to_pptx.style_parser import (
    parse_transform, apply_transform_to_pos, apply_gradient_fill
)
from pptx import Presentation
from pptx.util import Emu


def test_parse_transform_translate():
    tx, ty, rot, sx, sy = parse_transform("translate(100,50)")
    assert tx == pytest.approx(100.0)
    assert ty == pytest.approx(50.0)
    assert rot == pytest.approx(0.0)


def test_parse_transform_rotate():
    tx, ty, rot, sx, sy = parse_transform("rotate(45)")
    assert rot == pytest.approx(45.0)


def test_parse_transform_scale():
    tx, ty, rot, sx, sy = parse_transform("scale(2,3)")
    assert sx == pytest.approx(2.0)
    assert sy == pytest.approx(3.0)


def test_parse_transform_matrix():
    tx, ty, rot, sx, sy = parse_transform("matrix(1,0,0,1,50,30)")
    assert tx == pytest.approx(50.0)
    assert ty == pytest.approx(30.0)


def test_apply_transform_to_pos():
    x, y, w, h = apply_transform_to_pos(100, 200, 300, 400, "translate(50,25)")
    assert x == pytest.approx(150.0)
    assert y == pytest.approx(225.0)
    assert w == pytest.approx(300.0)


def test_apply_gradient_fill_linear():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Emu(0), Emu(0), Emu(1000), Emu(500))
    stops = [("0", "#ff0000"), ("100%", "#0000ff")]
    apply_gradient_fill(shape, stops, angle_deg=0.0)
    A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    from pptx.oxml.ns import qn
    spPr = shape._element.find(qn("p:spPr"))
    assert spPr.find(f"{{{A}}}gradFill") is not None
