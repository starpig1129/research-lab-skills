import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN

from svg_to_pptx.converter import CoordSystem
from svg_to_pptx.text_converter import add_textbox
from svg_to_pptx.style_parser import compute_style

CS = CoordSystem(svg_w=1200.0, svg_h=675.0)


def _slide():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_textbox_creates_shape():
    slide = _slide()
    elem = etree.fromstring('<text x="600" y="50" font-size="24">Hello</text>')
    add_textbox(slide, elem, compute_style(elem, {}), CS)
    assert len(slide.shapes) == 1


def test_textbox_text_content():
    slide = _slide()
    elem = etree.fromstring('<text x="100" y="100">World</text>')
    add_textbox(slide, elem, compute_style(elem, {}), CS)
    assert slide.shapes[0].text_frame.paragraphs[0].runs[0].text == "World"


def test_textbox_center_alignment():
    slide = _slide()
    elem = etree.fromstring('<text x="600" y="100" text-anchor="middle">Centered</text>')
    add_textbox(slide, elem, compute_style(elem, {}), CS)
    para = slide.shapes[0].text_frame.paragraphs[0]
    assert para.alignment == PP_ALIGN.CENTER


def test_textbox_tspan_multiline():
    slide = _slide()
    elem = etree.fromstring(
        '<text x="100" y="100">'
        '<tspan x="100" dy="0">Line 1</tspan>'
        '<tspan x="100" dy="20">Line 2</tspan>'
        '</text>'
    )
    add_textbox(slide, elem, compute_style(elem, {}), CS)
    tf = slide.shapes[0].text_frame
    texts = [p.runs[0].text for p in tf.paragraphs if p.runs]
    assert "Line 1" in texts
    assert "Line 2" in texts


def test_textbox_font_size():
    slide = _slide()
    elem = etree.fromstring('<text x="100" y="100" font-size="32">Big</text>')
    add_textbox(slide, elem, compute_style(elem, {}), CS)
    run = slide.shapes[0].text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(32)
