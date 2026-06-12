#!/usr/bin/env python3
"""make_pptx_png.py — Convert slide*.svg → PNG → deck.pptx (works on all PowerPoint versions)."""
import io, sys
from pathlib import Path
import cairosvg
from pptx import Presentation
from pptx.util import Emu

SLIDE_W = 12_192_000  # 16:9 widescreen EMU (33.87 cm)
SLIDE_H =  6_858_000

def main(slides_dir: str, out_path: str):
    slide_dir = Path(slides_dir)
    svg_files = sorted(slide_dir.glob("slide*.svg"))
    if not svg_files:
        print(f"No slide*.svg in {slide_dir}"); sys.exit(1)

    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for svg in svg_files:
        png = cairosvg.svg2png(url=str(svg.resolve()), output_width=1920, output_height=1080)
        slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
        slide.shapes.add_picture(io.BytesIO(png), 0, 0, Emu(SLIDE_W), Emu(SLIDE_H))
        print(f"  + {svg.name}")

    prs.save(out_path)
    print(f"\n{len(svg_files)} slide(s) → {out_path}")

if __name__ == "__main__":
    import argparse
    ap = argparse.ArgumentParser()
    ap.add_argument("--slides", required=True)
    ap.add_argument("--out",    required=True)
    args = ap.parse_args()
    main(args.slides, args.out)
