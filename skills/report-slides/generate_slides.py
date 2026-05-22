#!/usr/bin/env python3
"""generate_slides.py — Research Report SVG Slide Generator (Path A)

Usage:
    python scripts/generate_slides.py --data slide_data.json --out ./output/
    python scripts/generate_slides.py --data slide_data.json --out ./output/ --slide 5

Slide types:
    title, bullet_list, bar_chart, table, metric_cards, two_column, timeline, conclusion
"""

import json
import argparse
import os


# ── Style constants ───────────────────────────────────────────────────────────

S = {
    "w":      1200,
    "h":      675,
    "bg":     "#ffffff",
    "accent": "#1e3a5f",
    "good":   "#059669",
    "warn":   "#d97706",
    "danger": "#dc2626",
    "blue":   "#3b82f6",
    "card":   "#f8fafc",
    "border": "#e2e8f0",
    "body":   "#374151",
    "muted":  "#64748b",
    "white":  "#ffffff",
    "font":   "'Helvetica Neue', Arial, sans-serif",
    "top_bar_h": 6,
}

# ── Style loading ─────────────────────────────────────────────────────────────

def _parse_frontmatter(path: str) -> dict:
    """Parse YAML frontmatter from a .md file without requiring PyYAML."""
    try:
        with open(path, encoding="utf-8") as f:
            lines = f.readlines()
    except OSError as e:
        print(f"  [style] Cannot read {path}: {e}")
        return {}
    if not lines or lines[0].strip() != "---":
        return {}
    fm: dict = {}
    for line in lines[1:]:
        stripped = line.rstrip("\n")
        if stripped.strip() == "---":
            break
        if ":" not in stripped:
            continue
        key, _, val = stripped.partition(":")
        key = key.strip()
        val = val.strip().strip('"').strip("'")
        if key and val:
            fm[key] = val
    return fm


def apply_style(style_path: str) -> None:
    """Load a style .md file and override the global S dict."""
    fm = _parse_frontmatter(style_path)
    if not fm:
        return
    key_map = {
        "primary":  "accent",
        "bg":       "bg",
        "body":     "body",
        "muted":    "muted",
        "border":   "border",
        "card":     "card",
        "positive": "good",
        "warn":     "warn",
        "danger":   "danger",
        "font":     "font",
    }
    for style_key, s_key in key_map.items():
        if style_key in fm:
            S[s_key] = fm[style_key]
    if "top_bar_h" in fm:
        try:
            S["top_bar_h"] = int(fm["top_bar_h"])
        except ValueError:
            pass
    print(f"  [style] Applied: {style_path}")


# Chart drawing area
CL, CR, CT, CB = 130, 1100, 100, 520
CW, CH = CR - CL, CB - CT


# ── Text helpers ──────────────────────────────────────────────────────────────

def esc(v) -> str:
    return str(v).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace('"', "&quot;")


def wrap(text: str, max_chars: int = 60) -> list:
    words = str(text).split()
    lines, cur, n = [], [], 0
    for w in words:
        if n + len(w) + 1 > max_chars and cur:
            lines.append(" ".join(cur))
            cur, n = [w], len(w)
        else:
            cur.append(w)
            n += len(w) + 1
    if cur:
        lines.append(" ".join(cur))
    return lines or [""]


def tlines(lines: list, x, y, size, color,
           anchor="start", weight="normal", lh=1.45) -> str:
    spans = []
    for i, line in enumerate(lines):
        dy = "0" if i == 0 else f"{size * lh:.1f}"
        spans.append(f'<tspan x="{x}" dy="{dy}">{esc(line)}</tspan>')
    return (f'<text x="{x}" y="{y}" font-size="{size}" font-weight="{weight}" '
            f'fill="{color}" text-anchor="{anchor}">{"".join(spans)}</text>')


# ── Common slide frame ────────────────────────────────────────────────────────

def frame(title: str, footer: str = "") -> str:
    h = S["top_bar_h"]
    parts = [
        f'<rect width="{S["w"]}" height="{S["h"]}" fill="{S["bg"]}"/>',
        f'<rect x="0" y="0" width="{S["w"]}" height="{h}" fill="{S["accent"]}"/>',
        f'<text x="600" y="44" font-size="20" font-weight="700" fill="{S["accent"]}" '
        f'text-anchor="middle">{esc(title)}</text>',
        f'<line x1="40" y1="54" x2="1160" y2="54" stroke="{S["border"]}" stroke-width="1.5"/>',
    ]
    if footer:
        parts.append(f'<text x="1160" y="660" font-size="10" fill="{S["muted"]}" '
                     f'text-anchor="end">{esc(footer)}</text>')
    return "\n  ".join(parts)


def svg(body: str) -> str:
    return (f'<svg xmlns="http://www.w3.org/2000/svg" '
            f'viewBox="0 0 {S["w"]} {S["h"]}" font-family="{S["font"]}">\n'
            f'  {body}\n</svg>\n')


# ── Renderers ─────────────────────────────────────────────────────────────────

def render_title(sl: dict, meta: dict) -> str:
    title    = sl.get("title",    meta.get("experiment", "Research Report"))
    subtitle = sl.get("subtitle", "")
    author   = sl.get("author",   "")
    date     = sl.get("date",     meta.get("date", ""))
    footer   = meta.get("footer", "")
    cx       = 600

    parts = [
        f'<rect width="{S["w"]}" height="{S["h"]}" fill="{S["bg"]}"/>',
        f'<rect x="0" y="0" width="{S["w"]}" height="8" fill="{S["accent"]}"/>',
        f'<rect x="0" y="667" width="{S["w"]}" height="8" fill="{S["accent"]}"/>',
    ]

    title_lines = wrap(title, 52)
    title_y     = 255 - (len(title_lines) - 1) * 20
    parts.append(tlines(title_lines, cx, title_y, 30, S["accent"], "middle", "700", 1.3))

    div_y = title_y + len(title_lines) * 39
    parts.append(f'<line x1="200" y1="{div_y}" x2="1000" y2="{div_y}" '
                 f'stroke="{S["border"]}" stroke-width="1.5"/>')

    base_y = div_y + 35
    if subtitle:
        sub_lines = wrap(subtitle, 70)
        parts.append(tlines(sub_lines, cx, base_y, 16, S["muted"], "middle"))
        base_y += len(sub_lines) * 25 + 10

    meta_str = "  ·  ".join(filter(None, [author, date]))
    if meta_str:
        parts.append(f'<text x="{cx}" y="{base_y + 20}" font-size="13" '
                     f'fill="{S["muted"]}" text-anchor="middle">{esc(meta_str)}</text>')
    if footer:
        parts.append(f'<text x="1160" y="660" font-size="10" fill="{S["muted"]}" '
                     f'text-anchor="end">{esc(footer)}</text>')

    return svg("\n  ".join(parts))


def render_bullet_list(sl: dict, meta: dict) -> str:
    title    = sl.get("title", "")
    bullets  = sl.get("bullets", [])
    numbered = sl.get("numbered", False)
    footer   = meta.get("footer", "")

    parts = [frame(title, footer)]
    x_dot, x_text = 88, 116
    y = 90
    for i, item in enumerate(bullets):
        lines = wrap(str(item), 88)
        if numbered:
            parts.append(f'<circle cx="{x_dot}" cy="{y + 5}" r="12" fill="{S["accent"]}"/>')
            parts.append(f'<text x="{x_dot}" y="{y + 10}" font-size="11" font-weight="700" '
                         f'fill="{S["white"]}" text-anchor="middle">{i + 1}</text>')
        else:
            parts.append(f'<circle cx="{x_dot}" cy="{y + 6}" r="5" fill="{S["accent"]}"/>')
        parts.append(tlines(lines, x_text, y, 14, S["body"]))
        y += max(48, 24 * len(lines) + 16)

    return svg("\n  ".join(parts))


def render_bar_chart(sl: dict, meta: dict) -> str:
    title      = sl.get("title", "")
    categories = sl.get("categories", [])
    series     = sl.get("series", [])
    y_max      = float(sl.get("y_max", 100))
    note       = sl.get("note", "")
    footer     = meta.get("footer", "")

    parts = [frame(title, footer)]

    # Gridlines & Y-axis labels
    for i in range(6):
        val = y_max * i / 5
        y   = CB - (val / y_max) * CH
        parts.append(f'<line x1="{CL}" y1="{y:.1f}" x2="{CR}" y2="{y:.1f}" '
                     f'stroke="{S["border"]}" stroke-width="1"/>')
        parts.append(f'<text x="{CL - 8}" y="{y + 4:.1f}" font-size="10" '
                     f'fill="{S["muted"]}" text-anchor="end">{val:.0f}%</text>')

    parts.append(f'<line x1="{CL}" y1="{CT}" x2="{CL}" y2="{CB}" '
                 f'stroke="{S["muted"]}" stroke-width="1.5"/>')
    parts.append(f'<line x1="{CL}" y1="{CB}" x2="{CR}" y2="{CB}" '
                 f'stroke="{S["muted"]}" stroke-width="1.5"/>')

    n_cats = len(categories)
    n_ser  = len(series)
    if not n_cats or not n_ser:
        return svg("\n  ".join(parts))

    cat_slot = CW / n_cats
    group_w  = cat_slot * 0.70
    bar_w    = group_w / n_ser
    pad      = cat_slot * 0.15

    for ci, cat in enumerate(categories):
        gx = CL + ci * cat_slot + pad
        lx = gx + group_w / 2
        parts.append(f'<text x="{lx:.1f}" y="{CB + 20}" font-size="12" font-weight="600" '
                     f'fill="{S["body"]}" text-anchor="middle">{esc(cat)}</text>')

        for si, ser in enumerate(series):
            vals  = ser.get("values", [])
            if ci >= len(vals):
                continue
            val   = float(vals[ci])
            color = ser.get("color", S["blue"])
            bx    = gx + si * bar_w
            bh    = max((val / y_max) * CH, 2)
            by    = CB - bh

            parts.append(f'<rect x="{bx:.1f}" y="{by:.1f}" '
                         f'width="{bar_w - 3:.1f}" height="{bh:.1f}" fill="{color}"/>')
            if bh > 16:
                parts.append(f'<text x="{bx + (bar_w - 3) / 2:.1f}" y="{by - 4:.1f}" '
                             f'font-size="11" font-weight="700" fill="{color}" '
                             f'text-anchor="middle">{val:.1f}%</text>')

    # Legend
    lx = CL
    for si, ser in enumerate(series):
        color = ser.get("color", S["blue"])
        parts.append(f'<rect x="{lx + si * 230}" y="{CB + 40}" '
                     f'width="16" height="12" fill="{color}"/>')
        parts.append(f'<text x="{lx + si * 230 + 22}" y="{CB + 51}" '
                     f'font-size="12" fill="{S["body"]}">{esc(ser.get("label", ""))}</text>')

    if note:
        parts.append(f'<text x="{CR}" y="{CB + 51}" font-size="10" '
                     f'fill="{S["muted"]}" text-anchor="end">{esc(note)}</text>')

    return svg("\n  ".join(parts))


def render_table(sl: dict, meta: dict) -> str:
    title         = sl.get("title", "")
    columns       = sl.get("columns", [])
    rows          = sl.get("rows", [])
    highlight_col = sl.get("highlight_col")   # 0-indexed; colorizes +/- values
    footer        = meta.get("footer", "")

    parts = [frame(title, footer)]

    n_cols = len(columns)
    if not n_cols:
        return svg("\n  ".join(parts))

    tl, tr = 60, 1140
    tw     = tr - tl
    col_w  = tw / n_cols
    row_h  = min(50, 450 / (len(rows) + 1))
    top_y  = 75

    # Header
    parts.append(f'<rect x="{tl}" y="{top_y}" width="{tw}" '
                 f'height="{row_h}" fill="{S["accent"]}"/>')
    for ci, col in enumerate(columns):
        cx = tl + ci * col_w + col_w / 2
        parts.append(f'<text x="{cx:.1f}" y="{top_y + row_h * 0.63:.1f}" '
                     f'font-size="13" font-weight="700" fill="{S["white"]}" '
                     f'text-anchor="middle">{esc(col)}</text>')

    # Data rows
    for ri, row in enumerate(rows):
        ry = top_y + (ri + 1) * row_h
        bg = S["card"] if ri % 2 == 0 else S["bg"]
        parts.append(f'<rect x="{tl}" y="{ry:.1f}" width="{tw}" '
                     f'height="{row_h:.1f}" fill="{bg}" '
                     f'stroke="{S["border"]}" stroke-width="0.5"/>')
        for ci, cell in enumerate(row):
            cx    = tl + ci * col_w + col_w / 2
            cy    = ry + row_h * 0.63
            color = S["body"]
            if highlight_col is not None and ci == highlight_col:
                cs = str(cell)
                if "+" in cs:
                    color = S["good"]
                elif "-" in cs:
                    color = S["danger"]
            parts.append(f'<text x="{cx:.1f}" y="{cy:.1f}" font-size="13" '
                         f'fill="{color}" text-anchor="middle">{esc(cell)}</text>')

    parts.append(f'<rect x="{tl}" y="{top_y}" width="{tw}" '
                 f'height="{row_h * (len(rows) + 1):.1f}" fill="none" '
                 f'stroke="{S["border"]}" stroke-width="1.5"/>')

    return svg("\n  ".join(parts))


def render_metric_cards(sl: dict, meta: dict) -> str:
    title   = sl.get("title", "")
    metrics = sl.get("metrics", [])
    footer  = meta.get("footer", "")

    parts = [frame(title, footer)]
    n = len(metrics)
    if not n:
        return svg("\n  ".join(parts))

    cols = 2 if n == 4 else min(n, 3)
    rows = (n + cols - 1) // cols
    pad  = 40
    gap  = 20
    cw   = (S["w"] - 2 * pad - (cols - 1) * gap) / cols
    ch   = (S["h"] - 80 - 2 * pad - (rows - 1) * gap) / rows

    for i, m in enumerate(metrics):
        col   = i % cols
        row   = i // cols
        cx    = pad + col * (cw + gap)
        cy    = 80 + pad + row * (ch + gap)
        color = m.get("color", S["blue"])
        label = m.get("label", "")
        value = m.get("value", "")
        change = m.get("change", "")

        parts.append(f'<rect x="{cx:.1f}" y="{cy:.1f}" width="{cw:.1f}" height="{ch:.1f}" '
                     f'rx="8" fill="{S["card"]}" stroke="{S["border"]}" stroke-width="1.5"/>')
        parts.append(f'<rect x="{cx:.1f}" y="{cy:.1f}" width="{cw:.1f}" height="5" '
                     f'rx="4" fill="{color}"/>')
        parts.append(f'<text x="{cx + cw/2:.1f}" y="{cy + 36:.1f}" font-size="13" '
                     f'fill="{S["muted"]}" text-anchor="middle">{esc(label)}</text>')
        parts.append(f'<text x="{cx + cw/2:.1f}" y="{cy + ch/2 + 16:.1f}" font-size="38" '
                     f'font-weight="700" fill="{color}" text-anchor="middle">{esc(value)}</text>')
        if change:
            cc = S["good"] if "+" in str(change) else (S["danger"] if "-" in str(change) else S["muted"])
            parts.append(f'<text x="{cx + cw/2:.1f}" y="{cy + ch - 18:.1f}" '
                         f'font-size="12" fill="{cc}" text-anchor="middle">{esc(change)}</text>')

    return svg("\n  ".join(parts))


def render_two_column(sl: dict, meta: dict) -> str:
    title  = sl.get("title", "")
    left   = sl.get("left",  {})
    right  = sl.get("right", {})
    footer = meta.get("footer", "")

    parts = [frame(title, footer)]

    def panel(p: dict, px, py, pw, ph) -> list:
        out = []
        out.append(f'<rect x="{px}" y="{py}" width="{pw}" height="{ph}" '
                   f'rx="6" fill="{S["card"]}" stroke="{S["border"]}" stroke-width="1.5"/>')
        pt = p.get("title", "")
        if pt:
            out.append(f'<text x="{px + 20}" y="{py + 30}" font-size="14" font-weight="700" '
                       f'fill="{S["accent"]}">{esc(pt)}</text>')
            out.append(f'<line x1="{px + 20}" y1="{py + 38}" x2="{px + pw - 20}" y2="{py + 38}" '
                       f'stroke="{S["border"]}" stroke-width="1"/>')

        content = p.get("content", [])
        ty = py + 60
        max_c = int(pw / 9)

        if isinstance(content, str):
            out.append(tlines(wrap(content, max_c), px + 20, ty, 13, S["body"]))
        else:
            for item in content:
                lines = wrap(str(item), max_c - 4)
                out.append(f'<circle cx="{px + 28}" cy="{ty + 3}" r="4" fill="{S["accent"]}"/>')
                out.append(tlines(lines, px + 44, ty, 13, S["body"]))
                ty += 24 * len(lines) + 8
        return out

    parts += panel(left,  60,  66, 500, 562)
    parts += panel(right, 640, 66, 500, 562)

    return svg("\n  ".join(parts))


def render_timeline(sl: dict, meta: dict) -> str:
    title  = sl.get("title", "")
    events = sl.get("events", [])
    footer = meta.get("footer", "")

    parts = [frame(title, footer)]
    if not events:
        return svg("\n  ".join(parts))

    n  = len(events)
    y0 = 370
    x0, x1 = 100, 1100
    xs = [x0 + i * (x1 - x0) / max(n - 1, 1) for i in range(n)]

    parts.append(f'<line x1="{x0}" y1="{y0}" x2="{x1}" y2="{y0}" '
                 f'stroke="{S["border"]}" stroke-width="3"/>')

    for i, (ev, x) in enumerate(zip(events, xs)):
        color = ev.get("color", S["accent"])
        above = i % 2 == 0

        parts.append(f'<circle cx="{x:.1f}" cy="{y0}" r="10" '
                     f'fill="{color}" stroke="{S["bg"]}" stroke-width="2"/>')

        # Connector
        if above:
            parts.append(f'<line x1="{x:.1f}" y1="{y0 - 10}" x2="{x:.1f}" y2="{y0 - 28}" '
                         f'stroke="{color}" stroke-width="1.5" stroke-dasharray="3,2"/>')
            ty = y0 - 46
        else:
            parts.append(f'<line x1="{x:.1f}" y1="{y0 + 10}" x2="{x:.1f}" y2="{y0 + 28}" '
                         f'stroke="{color}" stroke-width="1.5" stroke-dasharray="3,2"/>')
            ty = y0 + 42

        label = ev.get("label", "")
        label_lines = wrap(label, 18)
        parts.append(tlines(label_lines, x, ty, 13, S["accent"], "middle", "700"))

        date_str = ev.get("date", "")
        if date_str:
            dy = y0 + 28 if above else y0 - 20
            parts.append(f'<text x="{x:.1f}" y="{dy}" font-size="10" '
                         f'fill="{S["muted"]}" text-anchor="middle">{esc(date_str)}</text>')

        detail = ev.get("detail", "")
        if detail:
            det_y = ty + len(label_lines) * 18 + 6
            parts.append(tlines(wrap(detail, 20), x, det_y, 11, S["muted"], "middle"))

    return svg("\n  ".join(parts))


def render_conclusion(sl: dict, meta: dict) -> str:
    title       = sl.get("title", "Conclusion & Next Steps")
    conclusions = sl.get("conclusions", [])
    next_steps  = sl.get("next_steps", [])
    footer      = meta.get("footer", "")
    c_head      = sl.get("conclusions_heading", "Conclusions")
    n_head      = sl.get("next_steps_heading",  "Next Steps")

    parts = [frame(title, footer)]

    def block(items, px, py, pw, ph, color, heading, numbered=True) -> list:
        out = []
        out.append(f'<rect x="{px}" y="{py}" width="{pw}" height="{ph}" '
                   f'rx="6" fill="{S["card"]}" stroke="{S["border"]}" stroke-width="1.5"/>')
        out.append(f'<rect x="{px}" y="{py}" width="{pw}" height="5" rx="4" fill="{color}"/>')
        out.append(f'<text x="{px + 20}" y="{py + 32}" font-size="14" font-weight="700" '
                   f'fill="{color}">{esc(heading)}</text>')
        out.append(f'<line x1="{px + 20}" y1="{py + 40}" x2="{px + pw - 20}" y2="{py + 40}" '
                   f'stroke="{S["border"]}" stroke-width="1"/>')

        iy = py + 62
        mc = int(pw / 9)
        for idx, item in enumerate(items):
            lines = wrap(str(item), mc - 4)
            if numbered:
                out.append(f'<circle cx="{px + 28}" cy="{iy + 2}" r="11" fill="{color}"/>')
                out.append(f'<text x="{px + 28}" y="{iy + 7}" font-size="10" font-weight="700" '
                           f'fill="{S["white"]}" text-anchor="middle">{idx + 1}</text>')
                out.append(tlines(lines, px + 50, iy, 13, S["body"]))
            else:
                out.append(f'<circle cx="{px + 28}" cy="{iy + 4}" r="5" fill="{color}"/>')
                out.append(tlines(lines, px + 46, iy, 13, S["body"]))
            iy += 28 * len(lines) + 10
        return out

    parts += block(conclusions, 60,  66, 500, 562, S["accent"], c_head, numbered=False)
    parts += block(next_steps,  640, 66, 500, 562, S["good"],   n_head, numbered=True)

    return svg("\n  ".join(parts))


# ── Dispatch ──────────────────────────────────────────────────────────────────

RENDERERS = {
    "title":         render_title,
    "bullet_list":   render_bullet_list,
    "bar_chart":     render_bar_chart,
    "table":         render_table,
    "metric_cards":  render_metric_cards,
    "two_column":    render_two_column,
    "timeline":      render_timeline,
    "conclusion":    render_conclusion,
}


def generate_slide(sl: dict, meta: dict) -> str:
    renderer = RENDERERS.get(sl["type"])
    if not renderer:
        raise ValueError(
            f"Unknown type: {sl['type']!r}. Available: {list(RENDERERS)}")
    return renderer(sl, meta)


# ── SVG → PNG conversion (fallback chain) ────────────────────────────────────

def svg_to_png(svg_path: str, png_path: str, width: int = 2400, height: int = 1350) -> bool:
    """Convert one SVG to PNG. Tries cairosvg → inkscape → rsvg-convert."""
    # 1. cairosvg (pip install cairosvg)
    try:
        import cairosvg
        cairosvg.svg2png(url=svg_path, write_to=png_path,
                         output_width=width, output_height=height)
        return True
    except ImportError:
        pass
    except Exception as e:
        print(f"    cairosvg error: {e}")

    import subprocess, shutil

    # 2. inkscape (system package)
    if shutil.which("inkscape"):
        r = subprocess.run(
            ["inkscape", svg_path,
             f"--export-filename={png_path}",
             f"--export-width={width}", f"--export-height={height}"],
            capture_output=True,
        )
        if r.returncode == 0:
            return True

    # 3. rsvg-convert (librsvg)
    if shutil.which("rsvg-convert"):
        r = subprocess.run(
            ["rsvg-convert", svg_path,
             "-w", str(width), "-h", str(height), "-o", png_path],
            capture_output=True,
        )
        if r.returncode == 0:
            return True

    return False


# ── SVG dir → PPTX ───────────────────────────────────────────────────────────

def to_pptx(svg_dir: str, output_path: str) -> None:
    """Convert all slide*.svg in svg_dir (sorted) into a single PPTX file."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        print("Missing dependency: python-pptx")
        print("Install with: pip install python-pptx")
        return

    import glob, tempfile

    svgs = sorted(glob.glob(os.path.join(svg_dir, "slide*.svg")))
    if not svgs:
        print(f"No slide*.svg found in {svg_dir}")
        return

    # Check at least one PNG converter is available before starting
    import shutil
    has_converter = (
        _try_import("cairosvg")
        or shutil.which("inkscape")
        or shutil.which("rsvg-convert")
    )
    if not has_converter:
        print("No SVG→PNG converter found. Install one of:")
        print("  pip install cairosvg")
        print("  sudo apt install inkscape   (or librsvg2-bin for rsvg-convert)")
        return

    prs = Presentation()
    prs.slide_width  = Inches(13.333)   # 16:9, matches 1200px wide
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]        # blank layout

    print(f"Converting {len(svgs)} slides → {output_path}")

    with tempfile.TemporaryDirectory() as tmp:
        for svg_path in svgs:
            name     = os.path.basename(svg_path).replace(".svg", ".png")
            png_path = os.path.join(tmp, name)

            ok = svg_to_png(svg_path, png_path)
            if not ok:
                print(f"  ✗ {os.path.basename(svg_path)} (conversion failed, skipped)")
                continue

            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(
                png_path,
                left=0, top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )
            print(f"  ✓ {os.path.basename(svg_path)}")

    prs.save(output_path)
    size_kb = os.path.getsize(output_path) // 1024
    print(f"\n✅ PPTX saved: {output_path}  ({size_kb} KB, {len(prs.slides)} slides)")
    print("\n注意：PPTX 中每張投影片為嵌入圖片，元素不可在 PowerPoint 內單獨編輯。")
    print("若需編輯單張，請修改對應 SVG 再重新執行 --to-pptx。")


def _try_import(name: str) -> bool:
    try:
        __import__(name)
        return True
    except ImportError:
        return False


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    ap = argparse.ArgumentParser(description="Generate SVG research slides from JSON")
    ap.add_argument("--data",     help="Path to slide_data.json")
    ap.add_argument("--out",      help="Output directory for SVG files")
    ap.add_argument("--slide",    type=int, default=None, help="Render only slide N")
    ap.add_argument("--style",    metavar="FILE",      default=None,
                    help="Style .md file to override colors/fonts (see styles/STYLES.md)")
    ap.add_argument("--to-pptx",  metavar="SVG_DIR",
                    help="Convert all slide*.svg in SVG_DIR to PPTX (skips SVG generation)")
    ap.add_argument("--pptx-out", metavar="FILE", default=None,
                    help="Output PPTX path (default: SVG_DIR/deck.pptx)")
    args = ap.parse_args()

    # ── Mode: convert existing SVGs to PPTX ──
    if args.to_pptx:
        out = args.pptx_out or os.path.join(args.to_pptx, "deck.pptx")
        to_pptx(args.to_pptx, out)
        return

    # ── Mode: generate SVGs from JSON ──
    if not args.data or not args.out:
        ap.error("--data and --out are required when not using --to-pptx")

    if args.style:
        apply_style(args.style)

    with open(args.data, encoding="utf-8") as f:
        data = json.load(f)

    meta   = data.get("meta", {})
    slides = data.get("slides", [])
    os.makedirs(args.out, exist_ok=True)

    generated = 0
    for sl in slides:
        if args.slide is not None and sl.get("index") != args.slide:
            continue
        content  = generate_slide(sl, meta)
        idx      = sl.get("index", 0)
        filename = f"slide{idx:02d}_{sl['type']}.svg"
        path     = os.path.join(args.out, filename)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        print(f"  ✓ {filename}")
        generated += 1

    print(f"\n{generated} slide(s) written to {args.out}")


if __name__ == "__main__":
    main()
