"""converter.py — orchestration layer for SVG → native PPTX conversion."""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Any

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

PPTX_W = 12_192_000
PPTX_H = 6_858_000
_SVG_NS = "http://www.w3.org/2000/svg"
_XLINK_NS = "http://www.w3.org/1999/xlink"


@dataclass
class CoordSystem:
    svg_w: float
    svg_h: float

    def x(self, v: float) -> int:
        return round(float(v) / self.svg_w * PPTX_W)

    def y(self, v: float) -> int:
        return round(float(v) / self.svg_h * PPTX_H)

    @classmethod
    def from_viewbox(cls, viewbox: str) -> "CoordSystem":
        parts = viewbox.strip().split()
        if len(parts) == 4:
            try:
                return cls(svg_w=float(parts[2]), svg_h=float(parts[3]))
            except ValueError:
                pass
        return cls(svg_w=1200.0, svg_h=675.0)


def _local_tag(elem: Any) -> str:
    tag = elem.tag
    if not isinstance(tag, str):
        return ""
    return tag.split("}")[-1] if "}" in tag else tag


class SvgConverter:
    """Converts one SVG file into one PPTX slide."""

    def __init__(self, svg_path: str, verbose: bool = False) -> None:
        self.svg_path = svg_path
        self.verbose = verbose
        tree = etree.parse(svg_path)
        self.root = tree.getroot()
        vb = self.root.get("viewBox", "0 0 1200 675")
        self.cs = CoordSystem.from_viewbox(vb)
        self._defs: Dict[str, Any] = {}
        self._connector_registry: List = []
        self._shape_registry: List = []
        self._pending_texts: List = []

    def convert(self, prs: Presentation, slide_layout: Any) -> Any:
        self._connector_registry = []
        self._shape_registry = []
        self._pending_texts = []
        slide = prs.slides.add_slide(slide_layout)
        self._resolve_defs()
        self._resolve_use_elements()
        self._dispatch_children(slide, self.root, {})
        # Add all text boxes after shapes so they appear on top
        from .text_converter import add_textbox
        for text_elem, text_style in self._pending_texts:
            add_textbox(slide, text_elem, text_style, self.cs)
        self._bind_connectors(slide)
        return slide

    def _resolve_defs(self) -> None:
        for elem in self.root.iter():
            if _local_tag(elem) == "defs":
                for child in elem:
                    eid = child.get("id")
                    if eid:
                        self._defs[eid] = child

    def _resolve_use_elements(self) -> None:
        from copy import deepcopy
        for use_elem in list(self.root.iter()):
            if _local_tag(use_elem) != "use":
                continue
            href = use_elem.get("href") or use_elem.get(
                "{http://www.w3.org/1999/xlink}href", "")
            if not href.startswith("#"):
                continue
            ref_id = href[1:]
            ref = self._defs.get(ref_id)
            if ref is None:
                continue
            clone = deepcopy(ref)
            use_x = use_elem.get("x", "0")
            use_y = use_elem.get("y", "0")
            if use_x != "0" or use_y != "0":
                existing = clone.get("transform", "")
                clone.set("transform",
                          f"translate({use_x},{use_y}) {existing}".strip())
            parent = use_elem.getparent()
            if parent is not None:
                idx = list(parent).index(use_elem)
                parent.insert(idx, clone)
                parent.remove(use_elem)

    def _dispatch_children(self, slide: Any, parent: Any, inherited: Dict) -> None:
        from .style_parser import compute_style
        for elem in parent:
            tag = _local_tag(elem)
            try:
                style = compute_style(elem, inherited)
                self._dispatch_element(slide, elem, style)
            except Exception as exc:
                if self.verbose:
                    print(f"  [warn] {tag}: {exc}")

    def _dispatch_element(self, slide: Any, elem: Any, style: Dict) -> None:
        tag = _local_tag(elem)
        if tag in ("rect", "circle", "ellipse", "image"):
            from .shapes import dispatch_shape
            shape = dispatch_shape(slide, elem, style, self.cs, None)
            if shape is not None and tag in ("rect", "circle", "ellipse"):
                try:
                    if tag == "rect":
                        bx = float(elem.get("x", 0))
                        by = float(elem.get("y", 0))
                        bw = float(elem.get("width", 0))
                        bh = float(elem.get("height", 0))
                    elif tag == "circle":
                        cx = float(elem.get("cx", 0))
                        cy = float(elem.get("cy", 0))
                        r = float(elem.get("r", 0))
                        bx, by, bw, bh = cx - r, cy - r, 2 * r, 2 * r
                    else:
                        cx = float(elem.get("cx", 0))
                        cy = float(elem.get("cy", 0))
                        rx = float(elem.get("rx", 0))
                        ry = float(elem.get("ry", 0))
                        bx, by, bw, bh = cx - rx, cy - ry, 2 * rx, 2 * ry
                    self._shape_registry.append((elem, bx, by, bw, bh, shape.shape_id))
                except Exception:
                    pass
        elif tag == "text":
            self._pending_texts.append((elem, style))
        elif tag in ("line", "polyline", "polygon"):
            from .connector import dispatch_connector
            conns = dispatch_connector(slide, elem, style, self.cs)
            self._connector_registry.extend(
                [(conn, elem) for conn in conns]
            )
        elif tag == "path":
            from .path_parser import parse_path
            from .path_to_pptx import add_path_shape
            d = elem.get("d", "")
            if d:
                add_path_shape(slide, parse_path(d), self.cs, style)
        elif tag == "g":
            self._dispatch_children(slide, elem, style)

    def _bind_connectors(self, slide: Any) -> None:
        from .connector import build_anchor_map
        if not self._connector_registry or not self._shape_registry:
            return
        anchor_map = build_anchor_map(self._shape_registry)
        for conn, conn_elem in self._connector_registry:
            tag = _local_tag(conn_elem)
            if tag == "line":
                try:
                    begin_pt = (float(conn_elem.get("x1", 0)), float(conn_elem.get("y1", 0)))
                    end_pt = (float(conn_elem.get("x2", 0)), float(conn_elem.get("y2", 0)))
                    _try_bind(conn, begin_pt, end_pt, anchor_map)
                except Exception:
                    pass


def _try_bind(conn: Any, begin_pt: tuple, end_pt: tuple, anchor_map: Dict) -> None:
    import math
    from .connector import THRESHOLD, bind_connector_end
    for is_begin, pt in [(True, begin_pt), (False, end_pt)]:
        best_dist = float("inf")
        best_sp_id = None
        best_idx = None
        for sp_id, anchors in anchor_map.items():
            for ax, ay, aidx in anchors:
                d = math.hypot(pt[0] - ax, pt[1] - ay)
                if d < best_dist:
                    best_dist = d
                    best_sp_id = sp_id
                    best_idx = aidx
        if best_dist <= THRESHOLD and best_sp_id is not None:
            bind_connector_end(conn, is_begin, best_sp_id, best_idx)


def convert_file(slides_dir: str, out_path: str, verbose: bool = False) -> None:
    prs = Presentation()
    prs.slide_width = Emu(PPTX_W)
    prs.slide_height = Emu(PPTX_H)
    layout = prs.slide_layouts[6]

    svg_files = sorted(Path(slides_dir).glob("slide*.svg"))
    if not svg_files:
        raise ValueError(f"No slide*.svg files found in {slides_dir}")

    for svg_path in svg_files:
        conv = SvgConverter(str(svg_path), verbose=verbose)
        conv.convert(prs, layout)
        if verbose:
            print(f"  + {svg_path.name}")

    Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"\n{len(svg_files)} slide(s) → {out_path}")
