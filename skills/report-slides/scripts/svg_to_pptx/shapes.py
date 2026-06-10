"""shapes.py — rect, circle, ellipse, image → python-pptx shapes."""
from __future__ import annotations

import re
from typing import Any, Dict, Optional

from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN

from .style_parser import apply_fill, apply_stroke, resolve_color, compute_style
from .converter import CoordSystem, _local_tag

_ALIGN = {"start": PP_ALIGN.LEFT, "middle": PP_ALIGN.CENTER, "end": PP_ALIGN.RIGHT}


def dispatch_shape(slide: Any, elem: Any, style: Dict,
                   cs: CoordSystem, label_elem: Optional[Any]) -> Optional[Any]:
    tag = _local_tag(elem)
    if tag == "rect":
        return _add_rect(slide, elem, style, cs, label_elem)
    elif tag in ("circle", "ellipse"):
        return _add_oval(slide, elem, style, cs, label_elem)
    elif tag == "image":
        return _add_image(slide, elem, style, cs)
    return None


def _add_rect(slide: Any, elem: Any, style: Dict,
              cs: CoordSystem, label_elem: Optional[Any]) -> Any:
    x = cs.x(float(elem.get("x", 0)))
    y = cs.y(float(elem.get("y", 0)))
    w = max(1, cs.x(float(elem.get("width", 0))))
    h = max(1, cs.y(float(elem.get("height", 0))))
    shape = slide.shapes.add_shape(1, Emu(x), Emu(y), Emu(w), Emu(h))
    apply_fill(shape, style.get("fill", "black"))
    apply_stroke(shape, style)
    if label_elem is not None:
        _write_label(shape, label_elem, style)
    return shape


def _add_oval(slide: Any, elem: Any, style: Dict,
              cs: CoordSystem, label_elem: Optional[Any]) -> Any:
    tag = _local_tag(elem)
    if tag == "circle":
        cx = float(elem.get("cx", 0))
        cy = float(elem.get("cy", 0))
        r = float(elem.get("r", 0))
        x, y, w, h = cx - r, cy - r, 2 * r, 2 * r
    else:
        cx = float(elem.get("cx", 0))
        cy = float(elem.get("cy", 0))
        rx = float(elem.get("rx", 0))
        ry = float(elem.get("ry", 0))
        x, y, w, h = cx - rx, cy - ry, 2 * rx, 2 * ry
    ex = cs.x(x)
    ey = cs.y(y)
    ew = max(1, cs.x(w))
    eh = max(1, cs.y(h))
    shape = slide.shapes.add_shape(9, Emu(ex), Emu(ey), Emu(ew), Emu(eh))
    apply_fill(shape, style.get("fill", "black"))
    apply_stroke(shape, style)
    if label_elem is not None:
        _write_label(shape, label_elem, style)
    return shape


def _add_image(slide: Any, elem: Any, style: Dict, cs: CoordSystem) -> Optional[Any]:
    href = elem.get("href") or elem.get(
        "{http://www.w3.org/1999/xlink}href", "")
    if not href or href.startswith("data:"):
        return None
    x = cs.x(float(elem.get("x", 0)))
    y = cs.y(float(elem.get("y", 0)))
    w = max(1, cs.x(float(elem.get("width", 100))))
    h = max(1, cs.y(float(elem.get("height", 100))))
    try:
        return slide.shapes.add_picture(href, Emu(x), Emu(y), Emu(w), Emu(h))
    except Exception:
        return None


def _write_label(shape: Any, text_elem: Any, parent_style: Dict) -> None:
    tf = shape.text_frame
    tf.word_wrap = True
    lines = _collect_text_lines(text_elem, parent_style)
    for i, (text, line_style) in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        anchor = line_style.get("text-anchor", parent_style.get("text-anchor", "middle"))
        para.alignment = _ALIGN.get(anchor, PP_ALIGN.CENTER)
        run = para.add_run()
        run.text = text
        _apply_font(run, line_style, parent_style)


def _collect_text_lines(text_elem: Any, parent_style: Dict):
    lines = []
    direct_text = (text_elem.text or "").strip()
    if direct_text:
        lines.append((direct_text, parent_style))
    for tspan in text_elem:
        if _local_tag(tspan) == "tspan":
            ts = compute_style(tspan, parent_style)
            t = (tspan.text or "").strip()
            if t:
                lines.append((t, ts))
    if not lines:
        lines.append(("", parent_style))
    return lines


def _apply_font(run: Any, style: Dict, parent_style: Dict) -> None:
    size_raw = style.get("font-size", parent_style.get("font-size", "14"))
    try:
        run.font.size = Pt(float(re.sub(r"[^0-9.]", "", size_raw) or "14"))
    except ValueError:
        run.font.size = Pt(14)
    weight = style.get("font-weight", parent_style.get("font-weight", "normal"))
    if weight in ("bold", "700", "800", "900"):
        run.font.bold = True
    color_val = style.get("fill", parent_style.get("fill", "#000000"))
    rgb = resolve_color(color_val)
    if rgb:
        run.font.color.rgb = rgb
