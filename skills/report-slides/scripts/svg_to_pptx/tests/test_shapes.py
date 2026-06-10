import pytest
from lxml import etree
from pptx import Presentation
from pptx.util import Emu

from svg_to_pptx.converter import CoordSystem
from svg_to_pptx.shapes import dispatch_shape
from svg_to_pptx.style_parser import compute_style

CS = CoordSystem(svg_w=1200.0, svg_h=675.0)


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Emu(12_192_000)
    prs.slide_height = Emu(6_858_000)
    return prs.slides.add_slide(prs.slide_layouts[6]), prs


def test_rect_creates_one_shape():
    slide, _ = _blank_slide()
    elem = etree.fromstring('<rect x="100" y="50" width="200" height="100" fill="#3b82f6"/>')
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    assert len(slide.shapes) == 1


def test_rect_position_and_size():
    slide, _ = _blank_slide()
    elem = etree.fromstring('<rect x="0" y="0" width="1200" height="675" fill="#fff"/>')
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    shape = slide.shapes[0]
    assert shape.left == 0
    assert shape.top == 0
    assert shape.width == 12_192_000
    assert shape.height == 6_858_000


def test_circle_creates_oval():
    slide, _ = _blank_slide()
    elem = etree.fromstring('<circle cx="600" cy="337" r="100" fill="#10b981"/>')
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    assert len(slide.shapes) == 1
    shape = slide.shapes[0]
    expected_w = CS.x(200)
    assert abs(shape.width - expected_w) <= 1


def test_rect_with_label_writes_text_in_shape():
    slide, _ = _blank_slide()
    rect_elem = etree.fromstring('<rect x="40" y="80" width="160" height="70" fill="#3b82f6"/>')
    text_elem = etree.fromstring('<text x="120" y="115" fill="white">Hello</text>')
    style = compute_style(rect_elem, {})
    dispatch_shape(slide, rect_elem, style, CS, text_elem)
    shape = slide.shapes[0]
    assert shape.has_text_frame
    assert shape.text_frame.paragraphs[0].runs[0].text == "Hello"


def test_ellipse_creates_oval():
    slide, _ = _blank_slide()
    elem = etree.fromstring('<ellipse cx="300" cy="200" rx="80" ry="40" fill="#f00"/>')
    style = compute_style(elem, {})
    dispatch_shape(slide, elem, style, CS, None)
    assert len(slide.shapes) == 1
