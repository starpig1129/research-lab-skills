"""text_converter.py — standalone TextBox for SVG <text> elements."""
from __future__ import annotations

import re
from typing import Any, Dict

from pptx.util import Emu
from pptx.enum.text import PP_ALIGN

from .style_parser import compute_style
from .converter import CoordSystem, _local_tag
from .shapes import _apply_font, _ALIGN

_ASCENT_FACTOR = 0.75
_CHAR_W_FACTOR = 0.65   # approx char width as fraction of font-size for Helvetica/Arial
_MIN_TB_W_SVG = 30.0
_PAD_SVG = 24.0         # horizontal padding in SVG units


def _font_size_svg(style: Dict) -> float:
    raw = style.get("font-size", "14")
    try:
        return float(re.sub(r"[^0-9.]", "", raw) or "14")
    except ValueError:
        return 14.0


def _tspan_dy(ts: Any) -> float:
    try:
        return float(ts.get("dy", "0") or "0")
    except (ValueError, TypeError):
        return 0.0


def _estimate_text_width(elem: Any, fs: float) -> float:
    """Estimate textbox width in SVG units based on longest text line."""
    max_chars = 0
    dt = (elem.text or "").strip()
    if dt:
        max_chars = max(max_chars, len(dt))
    for ts in elem:
        if _local_tag(ts) == "tspan":
            t = (ts.text or "").strip()
            if t:
                max_chars = max(max_chars, len(t))
    if max_chars == 0:
        max_chars = 6
    return max(max_chars * fs * _CHAR_W_FACTOR + _PAD_SVG, _MIN_TB_W_SVG)


def add_textbox(slide: Any, elem: Any, style: Dict, cs: CoordSystem) -> Any:
    try:
        tx = float(elem.get("x", 0))
        ty = float(elem.get("y", 0))
    except ValueError:
        tx, ty = 0.0, 0.0

    fs = _font_size_svg(style)
    anchor = style.get("text-anchor", "start")

    # SVG y is text baseline; PPTX textbox is positioned from top
    ty_top = ty - fs * _ASCENT_FACTOR

    tw = _estimate_text_width(elem, fs)
    if anchor == "middle":
        tx_left = tx - tw / 2
    elif anchor == "end":
        tx_left = tx - tw
    else:
        tx_left = tx

    tspans = [t for t in elem if _local_tag(t) == "tspan"]
    dy_values = [d for ts in tspans if (d := _tspan_dy(ts)) > 0]
    n_newlines = len(dy_values)
    dy_rep = dy_values[0] if dy_values else fs * 1.4
    tb_h_svg = fs * _ASCENT_FACTOR + n_newlines * dy_rep + fs * 0.3

    ex = cs.x(tx_left)
    ey = cs.y(ty_top)
    ew = cs.x(tw)
    eh = cs.y(max(tb_h_svg, fs * 1.3))

    tb = slide.shapes.add_textbox(Emu(ex), Emu(ey), Emu(ew), Emu(eh))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)

    _fill_text_frame(tf, elem, style, cs, anchor)
    return tb


def _fill_text_frame(tf: Any, elem: Any, style: Dict, cs: CoordSystem, anchor: str) -> None:
    tspans = [t for t in elem if _local_tag(t) == "tspan"]
    fs = _font_size_svg(style)
    align = _ALIGN.get(anchor, PP_ALIGN.LEFT)

    current_para = tf.paragraphs[0]
    current_para.alignment = align

    dt = (elem.text or "").strip()
    if dt:
        run = current_para.add_run()
        run.text = dt
        _apply_font(run, style, style)

    for ts in tspans:
        dy = _tspan_dy(ts)
        ts_style = compute_style(ts, style)
        ts_text = (ts.text or "").strip()
        ts_anchor = ts_style.get("text-anchor", anchor)
        ts_align = _ALIGN.get(ts_anchor, PP_ALIGN.LEFT)

        if dy > 0:
            p = tf.add_paragraph()
            p.alignment = ts_align
            space_svg = max(0.0, dy - fs)
            if space_svg > 0:
                p.space_before = Emu(cs.y(space_svg))
            current_para = p
        else:
            current_para.alignment = ts_align

        if ts_text:
            run = current_para.add_run()
            run.text = ts_text
            _apply_font(run, ts_style, style)
