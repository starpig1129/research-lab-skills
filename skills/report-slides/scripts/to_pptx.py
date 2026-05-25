#!/usr/bin/env python3
"""to_pptx.py — Pack SVG slides into a PPTX with native SVG embedding.

SVG files are embedded directly into the PPTX (viewed natively by PowerPoint 2016+/365).
A minimal white PNG is included as a fallback for older viewers.

Requirements:  pip install python-pptx
No cairosvg, Pillow, lxml, Inkscape, or any image converter needed.
Uses only the public python-pptx API + Python stdlib (zipfile, re, struct, zlib).
Works on Windows, macOS, and Linux.

Usage:
    python to_pptx.py --slides path/to/slides/ --out deck.pptx
"""

import argparse
import io
import re
import struct
import zipfile
import zlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

# ── Constants ─────────────────────────────────────────────────────────────────

SLIDE_W = 12_192_000   # EMU  (16:9 widescreen, 33.87 cm)
SLIDE_H =  6_858_000   # EMU

IMG_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
NS_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
NS_R    = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
SVG_EXT = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"


# ── Minimal white PNG — no Pillow required ────────────────────────────────────

def _make_white_png(w: int = 4, h: int = 3) -> bytes:
    """Generate a tiny white PNG using Python stdlib only."""
    def chunk(name: bytes, data: bytes) -> bytes:
        crc = zlib.crc32(name + data) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + name + data + struct.pack(">I", crc)

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)   # RGB 8-bit
    row  = b"\x00" + b"\xff\xff\xff" * w                   # filter byte + white pixels
    idat = zlib.compress(row * h)

    return (
        b"\x89PNG\r\n\x1a\n"
        + chunk(b"IHDR", ihdr)
        + chunk(b"IDAT", idat)
        + chunk(b"IEND", b"")
    )


_WHITE_PNG = _make_white_png()


# ── Core packing logic ────────────────────────────────────────────────────────

def pack_slides(svg_files: list, out_path: Path) -> None:
    """Build a PPTX from a list of SVG paths with native SVG embedding.

    Strategy:
      Step 1 — Build a base PPTX using only the public python-pptx API.
               Each slide gets a full-slide white PNG via add_picture().
      Step 2 — Post-process the PPTX ZIP with stdlib zipfile:
               • Add SVG files as media entries.
               • Register SVG relationships in each slide's .rels file.
               • Inject <p:extLst><asvg:svgBlip/> into each slide's <p:pic>.
               • Declare image/svg+xml in [Content_Types].xml.
    """

    # ── Step 1: public python-pptx API only ────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W)
    prs.slide_height = Emu(SLIDE_H)

    for _ in svg_files:
        slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank layout
        slide.shapes.add_picture(
            io.BytesIO(_WHITE_PNG), 0, 0, Emu(SLIDE_W), Emu(SLIDE_H)
        )

    raw = io.BytesIO()
    prs.save(raw)

    # ── Step 2: post-process ZIP ───────────────────────────────────────────────
    raw.seek(0)
    with zipfile.ZipFile(raw, "r") as zin:
        files = {name: zin.read(name) for name in zin.namelist()}

    # 2a. [Content_Types].xml — register image/svg+xml once
    ct_key = "[Content_Types].xml"
    if b"image/svg+xml" not in files[ct_key]:
        files[ct_key] = files[ct_key].replace(
            b"</Types>",
            b'<Default Extension="svg" ContentType="image/svg+xml"/></Types>',
        )

    # 2b. Per-slide: add relationship + media file
    slide_svg_rids: dict = {}

    for i, svg_path in enumerate(svg_files, start=1):
        rels_key  = f"ppt/slides/_rels/slide{i}.xml.rels"
        slide_key = f"ppt/slides/slide{i}.xml"

        if rels_key not in files or slide_key not in files:
            continue

        # Allocate next available rId
        existing_ids = [int(x) for x in re.findall(rb'Id="rId(\d+)"', files[rels_key])]
        svg_rid      = f"rId{max(existing_ids, default=0) + 1}"
        svg_target   = f"../media/slide{i:02d}_svg.svg"

        new_rel = (
            f'<Relationship Id="{svg_rid}" Type="{IMG_REL}" Target="{svg_target}"/>'
        ).encode()
        files[rels_key] = files[rels_key].replace(
            b"</Relationships>", new_rel + b"</Relationships>"
        )

        files[f"ppt/media/slide{i:02d}_svg.svg"] = Path(svg_path).read_bytes()
        slide_svg_rids[i] = svg_rid

    # 2c. Patch slide XMLs — append <p:extLst> inside the first <p:pic>
    for i, svg_rid in slide_svg_rids.items():
        slide_key = f"ppt/slides/slide{i}.xml"
        ext_block = (
            f'<p:extLst>'
            f'<p:ext uri="{SVG_EXT}">'
            f'<asvg:svgBlip'
            f' xmlns:asvg="{NS_ASVG}"'
            f' xmlns:r="{NS_R}"'
            f' r:embed="{svg_rid}"/>'
            f'</p:ext>'
            f'</p:extLst>'
        ).encode("utf-8")
        files[slide_key] = files[slide_key].replace(
            b"</p:pic>", ext_block + b"</p:pic>", 1
        )

    # ── Step 3: write patched ZIP to output path ───────────────────────────────
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(str(out_path), "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for name, data in files.items():
            zout.writestr(name, data)


# ── CLI ───────────────────────────────────────────────────────────────────────

def main() -> None:
    ap = argparse.ArgumentParser(
        description="Pack SVG slides into a PPTX with native SVG embedding"
    )
    ap.add_argument("--slides", required=True,
                    help="Directory containing slide*.svg files")
    ap.add_argument("--out",    required=True,
                    help="Output .pptx file path")
    args = ap.parse_args()

    slide_dir = Path(args.slides)
    svg_files = sorted(slide_dir.glob("slide*.svg"))

    if not svg_files:
        print(f"No slide*.svg files found in {slide_dir}")
        return

    out_path = Path(args.out)
    pack_slides(svg_files, out_path)

    for svg in svg_files:
        print(f"  + {svg.name}")
    print(f"\n{len(svg_files)} slide(s) -> {out_path}")
    print("Open with PowerPoint 2016+ / 365 for native SVG rendering.")
    print("Older viewers display the white PNG fallback (slide content is visible in SVG mode only).")


if __name__ == "__main__":
    main()
